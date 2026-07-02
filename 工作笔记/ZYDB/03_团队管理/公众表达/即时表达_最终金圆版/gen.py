# -*- coding: utf-8 -*-
"""金圆版生成 — python-pptx 在真实模板上操作 + LOGO 修复"""
import copy, io
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import sys; sys.stdout.reconfigure(encoding='utf-8')

TEMPLATE = r"D:\Data\04_DB\01_常态化工作\02_日志和汇报\02_汇报\02-4~5月经营分析会\Fw_ 关于召开公司月度经营分析会的通知\附件2：2026年4月及5月工作报告PPT模版二（非业务部门）.pptx"
OUTPUT = r"D:\Data\Note\工作笔记\ZYDB\03_团队管理\公众表达\风格对比\即时表达_最终金圆版.pptx"

# jinyuan_report.md 配色
P=RGBColor(0x1E,0x3A,0x5F); S=RGBColor(0x4A,0x90,0xD9); G=RGBColor(0xD4,0xAF,0x37)
BG=RGBColor(0xF5,0xF8,0xFC); T=RGBColor(0x33,0x33,0x33); TL=RGBColor(0x66,0x66,0x66); W=RGBColor(0xFF,0xFF,0xFF)

# 排版参数表（压缩版）
X1=Emu(300000); X2=Emu(4600000); CW=Emu(4200000); CY=Emu(850000); BAR_H=Emu(280000); FW=Emu(8400000)

prs = Presentation(TEMPLATE)

# 提取 LOGO
logo_blob=None
for sh in prs.slides[2].shapes:
    if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
        logo_blob=sh.image.blob
        lx,ly,lw,lh=sh.left,sh.top,sh.width,sh.height
        break

def dup(src):
    ns=prs.slides.add_slide(src.slide_layout)
    for s in list(ns.shapes): s._element.getparent().remove(s._element)
    for s in src.shapes: ns.shapes._spTree.append(copy.deepcopy(s.element))
    # 修复LOGO
    broken=[s for s in list(ns.shapes) if s.shape_type==MSO_SHAPE_TYPE.PICTURE]
    try:
        _=broken[0].image
    except:
        for b in broken: b._element.getparent().remove(b._element)
        if logo_blob: ns.shapes.add_picture(io.BytesIO(logo_blob),lx,ly,lw,lh)
    return ns

def clean(sl):
    rm=[]
    for s in sl.shapes:
        if s.name in("文本框 3","文本框 2","文本框 99"): rm.append(s)
        elif s.name=="标题 1": s.text_frame.clear()
    for s in rm: s._element.getparent().remove(s._element)

def find(sl,n):
    for s in sl.shapes:
        if s.name==n: return s

def rc(sl,x,y,w,h,f,st=None,rd=False):
    t=MSO_SHAPE.ROUNDED_RECTANGLE if rd else MSO_SHAPE.RECTANGLE
    s=sl.shapes.add_shape(t,x,y,w,h); s.fill.solid(); s.fill.fore_color.rgb=f
    if st: s.line.color.rgb=st; s.line.width=Pt(1)
    else: s.line.fill.background()
    return s

def ov(sl,x,y,w,h,f):
    s=sl.shapes.add_shape(MSO_SHAPE.OVAL,x,y,w,h); s.fill.solid(); s.fill.fore_color.rgb=f; s.line.fill.background(); return s

def tx(sl,x,y,w,h,t,sz=10,c=T,b=False,a=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=t; p.alignment=a
    for r in p.runs: r.font.size=Pt(sz); r.font.name="微软雅黑"; r.font.color.rgb=c; r.font.bold=b

def mt(sl,x,y,w,h,lines,sz=10,c=T):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,l in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=l
        for r in p.runs: r.font.size=Pt(sz); r.font.name="微软雅黑"; r.font.color.rgb=c

def tt(sl,text):
    s=find(sl,"文本框 26")
    if s:
        tf=s.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text="    "+text
        for r in p.runs: r.font.size=Pt(13); r.font.name="微软雅黑"; r.font.bold=True; r.font.color.rgb=P

tmpl=prs.slides[2]

# 封面
print("1.封面")
for sh in prs.slides[0].shapes:
    if sh.has_text_frame and("2026" in sh.text_frame.text or "汇报人" in sh.text_frame.text):
        tf=sh.text_frame; tf.clear()
        for i,(t,s,b) in enumerate([("即时表达",36,True),("公众发言的核心艺术",20,False),("",10,False),("汇报人：xxx       部门：xxx",12,False),("",10,False),("  2026年7月",12,False)]):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=t
            for r in p.runs: r.font.size=Pt(s); r.font.name="微软雅黑"; r.font.bold=b
        break

# P2 表达本质
print("2.表达本质")
s=dup(tmpl); clean(s); tt(s,"一、表达的本质：内容与方式")
rc(s,X1,CY,CW,Emu(2800000),BG,rd=True); rc(s,X1,CY,CW,BAR_H,S,rd=True)
tx(s,X1,CY+Emu(40000),CW,BAR_H,"表达内容 (WHAT)",12,W,True,PP_ALIGN.CENTER)
mt(s,X1+Emu(120000),CY+Emu(380000),CW-Emu(240000),Emu(2200000),["半书面 = 口头 + 底稿 + 审查","","● 前端：流畅自信的口头语言","● 后端：PPT/数据/逻辑链底稿","● 护城河：预判挑战点和答案","","经得起审查的公众发言"],9,T)
rc(s,X2,CY,CW,Emu(2800000),W,S,rd=True); rc(s,X2,CY,CW,BAR_H,P,rd=True)
tx(s,X2,CY+Emu(40000),CW,BAR_H,"表达方式 (HOW)",12,W,True,PP_ALIGN.CENTER)
mt(s,X2+Emu(120000),CY+Emu(380000),CW-Emu(240000),Emu(2200000),["本质：认知压榨","","① 流畅性启发","  感觉轻松 = 感觉正确","","② 场景表达模式","  钩子+结论+观点+行动","","更重要的是表达方式"],9,T)

# P3 四大原则
print("3.四大原则")
s=dup(tmpl); clean(s); tt(s,"二、表达内容四大原则")
for i,(n,t,d,c) in enumerate([("00","目的决定形式","没有野心，技巧沦为套路",G),("01","时间有限性","注意力有限，耐心有限",S),("02","时间线性","说完就没了，用短句口语词",P),("03","认知带宽","需要低压信息，废话是缓冲带",G)]):
    y=CY+i*Emu(700000); bc=P if i==3 else BG; tc=W if i==3 else T
    rc(s,X1,y,FW,Emu(620000),bc,rd=True); rc(s,X1,y,Emu(40000),Emu(620000),c)
    ov(s,X1+Emu(80000),y+Emu(130000),Emu(300000),Emu(300000),c)
    tx(s,X1+Emu(80000),y+Emu(130000),Emu(300000),Emu(300000),n,10,W,True,PP_ALIGN.CENTER)
    tx(s,X1+Emu(450000),y+Emu(100000),Emu(3000000),Emu(220000),t,13,tc,True)
    tx(s,X1+Emu(450000),y+Emu(350000),FW-Emu(550000),Emu(200000),d,10,TL if i!=3 else RGBColor(0xBD,0xD5,0xEA))

# P4 权威来源
print("4.权威来源")
s=dup(tmpl); clean(s); tt(s,"三、权威来源：信任 vs 权威")
rc(s,X1,CY,CW,Emu(1400000),RGBColor(0xFF,0xF8,0xE7),rd=True)
ov(s,X1+Emu(1100000),CY+Emu(150000),Emu(1800000),Emu(1800000),RGBColor(0xFF,0xFB,0xEB))
tx(s,X1+Emu(1100000),CY+Emu(600000),Emu(1800000),Emu(250000),"信任",20,P,True,PP_ALIGN.CENTER)
tx(s,X1+Emu(1100000),CY+Emu(900000),Emu(1800000),Emu(220000),"✓ 自信流畅  ✓ 眼神交流",10,T,False,PP_ALIGN.CENTER)
rc(s,X2,CY,CW,Emu(1400000),RGBColor(0xEA,0xF2,0xFB),rd=True)
ov(s,X2+Emu(1100000),CY+Emu(150000),Emu(1800000),Emu(1800000),RGBColor(0xEF,0xF6,0xFF))
tx(s,X2+Emu(1100000),CY+Emu(600000),Emu(1800000),Emu(250000),"权威",20,P,True,PP_ALIGN.CENTER)
tx(s,X2+Emu(1100000),CY+Emu(900000),Emu(1800000),Emu(220000),"✓ 数据支撑  ✓ 逻辑自洽",10,T,False,PP_ALIGN.CENTER)
rc(s,X1,CY+Emu(1550000),FW,Emu(500000),BG,rd=True)
tx(s,X1+Emu(80000),CY+Emu(1580000),FW-Emu(160000),Emu(220000),"信任：哪怕逻辑有瑕疵，直觉会被带着走",10,G,True,PP_ALIGN.CENTER)
tx(s,X1+Emu(80000),CY+Emu(1820000),FW-Emu(160000),Emu(220000),"权威：哪怕文笔枯燥，默认为专业资料",10,P,True,PP_ALIGN.CENTER)

# P5 流畅性启发
print("5.流畅性启发")
s=dup(tmpl); clean(s); tt(s,"四、流畅性启发（核心原理）")
cx=Emu(4072000)
ov(s,cx-Emu(550000),CY+Emu(50000),Emu(1100000),Emu(1100000),P)
ov(s,cx-Emu(430000),CY+Emu(170000),Emu(860000),Emu(860000),RGBColor(0x16,0x36,0x5A))
tx(s,cx-Emu(430000),CY+Emu(420000),Emu(860000),Emu(200000),"流畅性启发",14,W,True,PP_ALIGN.CENTER)
tx(s,cx-Emu(430000),CY+Emu(620000),Emu(860000),Emu(200000),"认知压榨本质",9,RGBColor(0x93,0xC5,0xFD),False,PP_ALIGN.CENTER)
rc(s,X1+Emu(900000),CY-Emu(100000),Emu(6600000),Emu(280000),BG,rd=True)
tx(s,X1+Emu(900000),CY-Emu(80000),Emu(6600000),Emu(280000),'大脑依赖"信息处理的主观体验"',10,T,False,PP_ALIGN.CENTER)
rc(s,X1,CY+Emu(1300000),Emu(3900000),Emu(300000),RGBColor(0xE8,0xF5,0xE9),rd=True)
tx(s,X1,CY+Emu(1320000),Emu(3900000),Emu(300000),"感觉轻松 = 感觉正确",11,RGBColor(0x2E,0x7D,0x32),True,PP_ALIGN.CENTER)
rc(s,X2+Emu(350000),CY+Emu(1300000),Emu(3900000),Emu(300000),RGBColor(0xFF,0xEB,0xEE),rd=True)
tx(s,X2+Emu(350000),CY+Emu(1320000),Emu(3900000),Emu(300000),"感觉困难 = 感觉可疑",11,RGBColor(0xC6,0x28,0x28),True,PP_ALIGN.CENTER)
rc(s,X1,CY+Emu(1800000),FW,Emu(350000),RGBColor(0xFF,0xF3,0xE0),rd=True)
tx(s,X1,CY+Emu(1830000),FW,Emu(350000),"启示：表达流畅比内容正确更重要",11,RGBColor(0xE8,0x77,0x22),True,PP_ALIGN.CENTER)

# P6 场景表达
print("6.场景表达")
s=dup(tmpl); clean(s); tt(s,"五、场景表达模式")
for i,(t,d,c,o) in enumerate([("钩子","关于X，我需要您",G,0),("结论","直接说核心结论",S,Emu(150000)),("观点摘要","重要观点概括",P,Emu(300000)),("行动指南","下一步建议",G,Emu(450000))]):
    y=CY+Emu(50000)+i*Emu(680000); x=X1+o; w=FW-o
    rc(s,x,y,w,Emu(580000),BG,rd=True); rc(s,x,y,Emu(40000),Emu(580000),c)
    ov(s,x+Emu(80000),y+Emu(120000),Emu(300000),Emu(300000),c)
    tx(s,x+Emu(80000),y+Emu(120000),Emu(300000),Emu(300000),str(i+1),10,W,True,PP_ALIGN.CENTER)
    tx(s,x+Emu(450000),y+Emu(100000),Emu(2500000),Emu(200000),t,13,P,True)
    tx(s,x+Emu(450000),y+Emu(330000),w-Emu(550000),Emu(200000),d,10,TL)

# P7 实战模板
print("7.实战模板")
s=dup(tmpl); clean(s); tt(s,"六、实战模板：内部汇报发言稿")
for i,(n,t,d) in enumerate([("1","结论放第一句","开门见时"),("2","2-3个理由","第一第二第三"),("3","建议收尾","我的建议是___")]):
    x=X1+i*Emu(2830000)
    rc(s,x,CY,Emu(2700000),Emu(1600000),BG,rd=True)
    ov(s,x+Emu(80000),CY+Emu(80000),Emu(300000),Emu(300000),G)
    tx(s,x+Emu(80000),CY+Emu(80000),Emu(300000),Emu(300000),n,10,W,True,PP_ALIGN.CENTER)
    tx(s,x+Emu(430000),CY+Emu(100000),Emu(2200000),Emu(200000),t,12,P,True)
    tx(s,x+Emu(80000),CY+Emu(500000),Emu(2540000),Emu(200000),d,10,T)
tx(s,X1,CY+Emu(1800000),FW,Emu(220000),"风格：自信果断 · 严谨冷静 · 口语化 · 短句",11,P,True)
rc(s,X1,CY+Emu(2100000),FW,Emu(300000),P,rd=True)
tx(s,X1,CY+Emu(2130000),FW,Emu(300000),'核心：让听众"轻松听懂"',11,G,True,PP_ALIGN.CENTER)

# 排序+结束页
print("8.排序")
sl=prs.slides._sldIdLst; al=list(sl)
for s in al: sl.remove(s)
for idx in [0,1,6,7,8,9,10,11,5]: sl.append(al[idx])
for sh in prs.slides[-1].shapes:
    if sh.has_text_frame and "汇报完毕" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs:
            for r in p.runs: r.text=r.text.replace("汇报完毕","感谢聆听")

prs.save(OUTPUT)
print(f"\n金圆版: {OUTPUT}")

# 验证LOGO
prs2=Presentation(OUTPUT)
for i,sl in enumerate(prs2.slides):
    for sh in sl.shapes:
        if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
            try: _=sh.image; print(f"S{i+1}: LOGO OK")
            except: print(f"S{i+1}: LOGO BROKEN")
